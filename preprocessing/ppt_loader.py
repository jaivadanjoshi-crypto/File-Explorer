# preprocessing/pptx_loader.py
from typing import List
from pptx import Presentation
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
import config

def load_and_split_pptx(uploaded_file) -> List[Document]:
    """Read a Streamlit PowerPoint file (.pptx) and extract text, ignoring images."""
    prs = Presentation(uploaded_file)
    slide_docs: List[Document] = []
    
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text_chunks = []
        
        for shape in slide.shapes:
            # Only extract from shapes containing text (skips raw image shapes automatically)
            if hasattr(shape, "text") and shape.text.strip():
                slide_text_chunks.append(shape.text.strip())
                
        text = "\n".join(slide_text_chunks)
        if not text.strip():
            continue
            
        slide_docs.append(
            Document(
                page_content=text,
                metadata={"source": uploaded_file.name, "slide": slide_num}
            )
        )
        
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=config.CHUNK_SIZE,
        chunk_overlap=config.CHUNK_OVERLAP,
    )
    return splitter.split_documents(slide_docs)
